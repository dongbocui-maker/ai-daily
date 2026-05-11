"""Low-level drawing helpers for FY27 IE-style slides."""

from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .style import C, FONT, SLIDE_H


def add_blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def set_bg(slide, color=None):
    if color is None:
        color = C["white"]
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    shp.shadow.inherit = False
    return shp


def add_textbox(slide, x, y, w, h, text, *,
                font_size=10, bold=False, italic=False, color=None,
                align="left", valign="top", char_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    if valign == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    if char_spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(char_spacing * 100)))
    return tb


def add_paragraphs(slide, x, y, w, h, lines, *,
                   font_size=10, color=None, bullet=False,
                   bullet_color=None, line_spacing=1.25,
                   bold=False, italic=False, align="left"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if not lines:
        return tb
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if bullet:
            br = p.add_run()
            br.text = "\u25CF  "
            br.font.name = FONT
            br.font.size = Pt(font_size)
            br.font.color.rgb = bullet_color or C["purple"]
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        if color is not None:
            run.font.color.rgb = color
    return tb


def add_header(slide, eyebrow=None, title=None, subtitle=None):
    """Standard FY27 IE page header (top zone, ~1.65 inch tall)."""
    add_textbox(slide, Inches(0.5), Inches(0.25), Inches(1.5), Inches(0.3),
                "accenture", font_size=11, bold=True, color=C["ink"], valign="middle")
    add_rect(slide, Inches(1.95), Inches(0.34), Inches(0.13), Inches(0.13), fill=C["purple"])
    if eyebrow:
        add_textbox(slide, Inches(0.5), Inches(0.72), Inches(12.3), Inches(0.3),
                    eyebrow.upper(), font_size=11, bold=True, color=C["purple"],
                    char_spacing=4)
    if title:
        add_textbox(slide, Inches(0.5), Inches(1.00), Inches(12.3), Inches(0.6),
                    title, font_size=26, bold=True, color=C["ink"])
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.58), Inches(12.3), Inches(0.32),
                    subtitle, font_size=11.5, italic=True, color=C["ink_soft"])


def add_footer(slide, page_no, total):
    add_textbox(slide, Inches(0.40), Inches(7.12), Inches(0.2), Inches(0.3),
                "\u203A", font_size=11, bold=True, color=C["purple"], valign="middle")
    add_textbox(slide, Inches(0.60), Inches(7.18), Inches(8), Inches(0.25),
                "Curated by aidigest.club  \u00b7  AI Daily Reads  \u00b7  Executive Brief",
                font_size=8, color=C["muted"])
    add_textbox(slide, Inches(11.5), Inches(7.18), Inches(1.3), Inches(0.25),
                f"{page_no} / {total}", font_size=8, color=C["muted"], align="right")
