#!/usr/bin/env python3
"""Generate FY27 IE Style PPTX for an outline JSON.

Usage:
    python generate_deck_fy27_ie.py <outline.json> [output.pptx]

Outline schema: see outlines/anthropic-81k-economics.json
"""
from __future__ import annotations
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from fy27_ie import style as S
from fy27_ie.builders_a import (
    build_cover, build_context, build_agenda, build_section_header,
)
from fy27_ie.builders_b import (
    build_stat_grid, build_three_card, build_two_column, build_quote_pair,
)
from fy27_ie.builders_c import (
    build_insight_pair, build_action_table, build_summary,
    build_references, build_closing,
)


BUILDERS = {
    "cover": build_cover,
    "context": build_context,
    "agenda": build_agenda,
    "section_header": build_section_header,
    "stat_grid": build_stat_grid,
    "three_card": build_three_card,
    "two_column": build_two_column,
    "quote_pair": build_quote_pair,
    "insight_pair": build_insight_pair,
    "action_table": build_action_table,
    "summary": build_summary,
    "references": build_references,
    "closing": build_closing,
}


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    outline_path = Path(sys.argv[1])
    if not outline_path.exists():
        print(f"outline file not found: {outline_path}")
        sys.exit(1)

    with open(outline_path, encoding="utf-8") as fh:
        outline = json.load(fh)

    slides = outline.get("slides", [])
    if not slides:
        print("outline has no slides")
        sys.exit(1)

    deck_meta = outline.get("deck_meta", {})
    out_path = Path(sys.argv[2]) if len(sys.argv) > 2 else (
        Path(__file__).parent / "output" / f"{outline.get('slug', 'deck')}-fy27ie.pptx"
    )
    out_path.parent.mkdir(parents=True, exist_ok=True)

    pres = Presentation()
    pres.slide_width = S.SLIDE_W
    pres.slide_height = S.SLIDE_H

    blank_layout = pres.slide_layouts[6]
    total = len(slides)

    for i, sd in enumerate(slides, 1):
        stype = sd.get("type", "cover")
        builder = BUILDERS.get(stype)
        sd = dict(sd)
        sd["_deck_meta"] = deck_meta
        slide = pres.slides.add_slide(blank_layout)
        if builder is None:
            print(f"  ! unknown slide type '{stype}' at index {i}, skipping")
            continue
        builder(slide, sd, i, total)
        print(f"  [{i:02d}/{total}] {stype:>16}  ok")

    pres.save(out_path)
    print(f"\n✅ saved: {out_path}")
    print(f"   size: {out_path.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
