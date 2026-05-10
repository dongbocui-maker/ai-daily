#!/usr/bin/env python3
"""
Generate Accenture-styled PPT deck from a structured outline JSON.

Usage:
    python3 generate_deck.py outlines/<slug>.json output/<slug>.pptx
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Brand tokens (Accenture-inspired) ----------
BRAND = {
    "purple_primary": RGBColor(0xA1, 0x00, 0xFF),
    "purple_dark":    RGBColor(0x46, 0x00, 0x73),
    "purple_mid":     RGBColor(0x7B, 0x42, 0xF6),
    "purple_light":   RGBColor(0xE5, 0xD4, 0xFF),
    "purple_tint":    RGBColor(0xF1, 0xE8, 0xFF),
    "berry":          RGBColor(0xBE, 0x82, 0xFF),
    "magenta":        RGBColor(0xFF, 0x50, 0xA0),
    "magenta_dark":   RGBColor(0xC8, 0x2A, 0x80),
    "ink":            RGBColor(0x14, 0x14, 0x14),
    "slate":          RGBColor(0x4A, 0x4A, 0x4A),
    "muted":          RGBColor(0x8A, 0x8A, 0x8A),
    "rule":           RGBColor(0xD9, 0xD9, 0xD9),
    "soft_rule":      RGBColor(0xEC, 0xEC, 0xEC),
    "bg_white":       RGBColor(0xFF, 0xFF, 0xFF),
    "bg_off":         RGBColor(0xFA, 0xF7, 0xFF),
    "evidence_bg":    RGBColor(0xF5, 0xEE, 0xFF),
}

FONT_HEAD = "Calibri"
FONT_BODY = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- Low-level helpers ----------

def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill, line=None, shape_type=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape_type, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def add_textbox(slide, left, top, width, height, text,
                font_name=FONT_BODY, size=18, bold=False,
                color=BRAND["ink"], align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, line_spacing=1.15, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, left, top, width, height, lines,
                   font_name=FONT_BODY, size=15, color=BRAND["ink"],
                   bullet=False, align=PP_ALIGN.LEFT, line_spacing=1.3,
                   space_after=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = (("• " + ln) if bullet else ln)
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_brand_mark(slide, left=Inches(0.4), top=Inches(0.4)):
    add_rect(slide, left, top, Inches(0.5), Inches(0.05),
             fill=BRAND["purple_primary"])


def add_kicker(slide, text, top=Inches(0.4)):
    add_textbox(slide, Inches(0.4), top, Inches(8), Inches(0.3),
                text.upper(), font_name=FONT_HEAD, size=10, bold=True,
                color=BRAND["purple_primary"])


def add_slide_title(slide, text, top=Inches(0.7)):
    add_textbox(slide, Inches(0.4), top, Inches(12.53), Inches(0.7),
                text, font_name=FONT_HEAD, size=28, bold=True,
                color=BRAND["purple_dark"])


def add_key_message(slide, text, top=Inches(1.4)):
    if not text:
        return
    add_textbox(slide, Inches(0.4), top, Inches(12.53), Inches(0.6),
                text, font_name=FONT_HEAD, size=16, bold=False,
                color=BRAND["slate"], line_spacing=1.25, italic=True)


def add_footer(slide, page_no, total, deck_meta):
    add_rect(slide, Inches(0.4), Inches(7.05), Inches(12.53), Emu(9525),
             fill=BRAND["soft_rule"])
    add_textbox(slide, Inches(0.4), Inches(7.12), Inches(7), Inches(0.3),
                f"Source: {deck_meta['source_authors']} · {deck_meta['source_paper']}",
                size=8, color=BRAND["muted"])
    add_textbox(slide, Inches(0.4), Inches(7.12), Inches(12.53), Inches(0.3),
                "Generative AI at Work · Executive Briefing",
                size=8, color=BRAND["muted"], align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.4), Inches(7.12), Inches(12.53), Inches(0.3),
                f"{page_no} / {total}",
                size=8, color=BRAND["muted"], align=PP_ALIGN.RIGHT)


def set_notes(slide, text):
    notes = slide.notes_slide
    nf = notes.notes_text_frame
    nf.text = text


# ---------- Slide builders ----------

def build_cover(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["bg_white"])

    # Big purple block + diagonal accent
    add_rect(slide, Inches(0), Inches(0), Inches(5.5), SLIDE_H, fill=BRAND["purple_dark"])
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, Inches(4.7), Inches(0), Inches(2.0), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND["purple_primary"]
    accent.line.fill.background()

    # Brand mark on dark block
    add_rect(slide, Inches(0.5), Inches(0.5), Inches(0.6), Inches(0.06), fill=BRAND["bg_white"])
    add_textbox(slide, Inches(0.5), Inches(0.65), Inches(4.5), Inches(0.4),
                "EXECUTIVE BRIEFING", size=11, bold=True, color=BRAND["bg_white"])
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(4.5), Inches(0.4),
                "AI Daily Reads · aidigest.club", size=10, color=BRAND["purple_light"])

    # Title block
    add_textbox(slide, Inches(7.0), Inches(1.6), Inches(6), Inches(1.5),
                data["title"], font_name=FONT_HEAD, size=44, bold=True,
                color=BRAND["purple_dark"], line_spacing=1.05)
    add_rect(slide, Inches(7.0), Inches(2.95), Inches(0.8), Inches(0.06),
             fill=BRAND["purple_primary"])
    add_textbox(slide, Inches(7.0), Inches(3.15), Inches(6), Inches(1.6),
                data["subtitle"], font_name=FONT_BODY, size=17,
                color=BRAND["slate"], line_spacing=1.3)

    # Highlights row (4 mini-stats)
    if "highlights" in data:
        hl_top = Inches(5.2)
        hl_h = Inches(1.0)
        hl_count = len(data["highlights"])
        hl_w_total = Inches(5.8)
        hl_w = Emu(int(hl_w_total) // hl_count - int(Inches(0.05)))
        for i, hl in enumerate(data["highlights"]):
            x = Inches(7.0) + Emu(int(hl_w_total) * i // hl_count)
            add_rect(slide, x, hl_top, hl_w, Emu(int(hl_h) - int(Inches(0.1))),
                     fill=BRAND["bg_off"], line=BRAND["rule"])
            add_textbox(slide, x + Inches(0.1), hl_top + Inches(0.1),
                        hl_w - Inches(0.2), Inches(0.4),
                        hl["value"], font_name=FONT_HEAD, size=14, bold=True,
                        color=BRAND["purple_primary"])
            add_textbox(slide, x + Inches(0.1), hl_top + Inches(0.45),
                        hl_w - Inches(0.2), Inches(0.4),
                        hl["label"], font_name=FONT_BODY, size=9,
                        color=BRAND["muted"])

    add_textbox(slide, Inches(7.0), Inches(6.7), Inches(6), Inches(0.4),
                data["footer_meta"], size=9, color=BRAND["muted"])

    set_notes(slide, "Cover slide — opens the briefing on the Brynjolfsson, Li & Raymond study of generative AI at work. Frame this as the first rigorous, large-scale, causal study of GenAI's impact on real workplaces (5,172 customer service agents). Key message before turning the page: GenAI is a skill-equalizing technology — fundamentally different from prior IT waves.")


def build_context(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["bg_white"])
    add_brand_mark(slide)
    add_kicker(slide, data["kicker"])
    add_slide_title(slide, data["title"])
    add_key_message(slide, data["key_message"])

    ba = data["before_after"]
    col_w = Inches(5.95)
    col_h = Inches(2.9)
    col_top = Inches(2.3)

    # BEFORE column
    add_rect(slide, Inches(0.4), col_top, col_w, Inches(0.5), fill=BRAND["muted"])
    add_textbox(slide, Inches(0.6), col_top + Inches(0.1), col_w - Inches(0.4), Inches(0.4),
                ba["before"]["header"].upper(), font_name=FONT_HEAD, size=12, bold=True,
                color=BRAND["bg_white"])
    add_rect(slide, Inches(0.4), col_top + Inches(0.5), col_w, col_h - Inches(0.5),
             fill=BRAND["bg_off"], line=BRAND["rule"])
    add_paragraphs(slide, Inches(0.7), col_top + Inches(0.7),
                   col_w - Inches(0.5), col_h - Inches(0.8),
                   ba["before"]["points"], size=14, color=BRAND["ink"],
                   bullet=True, line_spacing=1.4, space_after=4)

    # AFTER column
    add_rect(slide, Inches(6.95), col_top, col_w, Inches(0.5), fill=BRAND["purple_primary"])
    add_textbox(slide, Inches(7.15), col_top + Inches(0.1), col_w - Inches(0.4), Inches(0.4),
                ba["after"]["header"].upper(), font_name=FONT_HEAD, size=12, bold=True,
                color=BRAND["bg_white"])
    add_rect(slide, Inches(6.95), col_top + Inches(0.5), col_w, col_h - Inches(0.5),
             fill=BRAND["purple_tint"], line=BRAND["rule"])
    add_paragraphs(slide, Inches(7.25), col_top + Inches(0.7),
                   col_w - Inches(0.5), col_h - Inches(0.8),
                   ba["after"]["points"], size=14, color=BRAND["purple_dark"],
                   bullet=True, line_spacing=1.4, space_after=4)

    # Stat band at bottom
    band_top = col_top + col_h + Inches(0.25)
    band_h = Inches(0.85)
    sb = data["stat_band"]
    n = len(sb)
    band_w_total = Inches(12.53)
    cell_w = Emu(int(band_w_total) // n)
    for i, item in enumerate(sb):
        x = Inches(0.4) + Emu(int(cell_w) * i)
        add_rect(slide, x, band_top, cell_w - Inches(0.05), band_h,
                 fill=BRAND["purple_dark"])
        add_textbox(slide, x + Inches(0.2), band_top + Inches(0.08),
                    cell_w - Inches(0.4), Inches(0.4),
                    item["value"], font_name=FONT_HEAD, size=22, bold=True,
                    color=BRAND["bg_white"])
        add_textbox(slide, x + Inches(0.2), band_top + Inches(0.5),
                    cell_w - Inches(0.4), Inches(0.35),
                    item["label"], font_name=FONT_BODY, size=10,
                    color=BRAND["purple_light"])

    add_textbox(slide, Inches(0.4), band_top + band_h + Inches(0.1),
                Inches(12.53), Inches(0.25),
                f"Source reference: {data['source_note']}",
                size=8, color=BRAND["muted"])

    add_footer(slide, page_no, total, deck_meta)
    set_notes(slide, "Context slide. Set up why this paper is special: first causal, large-scale, real-workplace evidence on GenAI productivity. Walk left → right: previous evidence base (lab/surveys/anecdotes) versus what this paper adds (5,172 real agents, quasi-random rollout, DiD, outage-based natural experiment). The bottom band reinforces sample scale.")


def build_agenda(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["bg_off"])
    add_brand_mark(slide)
    add_kicker(slide, "Agenda")
    add_slide_title(slide, data["title"])

    items = data["items"]
    n = len(items)
    half = (n + 1) // 2
    item_h = Inches(0.95)
    start_top = Inches(2.0)
    col1_x = Inches(0.6)
    col2_x = Inches(7.0)

    for i, item in enumerate(items):
        if i < half:
            x, y = col1_x, start_top + item_h * i
        else:
            x, y = col2_x, start_top + item_h * (i - half)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BRAND["purple_primary"]
        circle.line.fill.background()
        add_textbox(slide, x, y + Inches(0.13), Inches(0.7), Inches(0.5),
                    f"{i+1:02d}", size=18, bold=True,
                    color=BRAND["bg_white"], align=PP_ALIGN.CENTER, font_name=FONT_HEAD)

        # Label + sub
        add_textbox(slide, x + Inches(0.9), y + Inches(0.05),
                    Inches(5.3), Inches(0.4),
                    item["label"] if isinstance(item, dict) else item,
                    font_name=FONT_HEAD, size=17, bold=True, color=BRAND["ink"])
        if isinstance(item, dict) and "sub" in item:
            add_textbox(slide, x + Inches(0.9), y + Inches(0.45),
                        Inches(5.3), Inches(0.5),
                        item["sub"], font_name=FONT_BODY, size=12,
                        color=BRAND["slate"], line_spacing=1.2)

    add_footer(slide, page_no, total, deck_meta)
    set_notes(slide, "Walk through the structure briefly. Pace: ~2 minutes per content section. Findings → Core Views → Themes → Insights → Actions → References.")


def build_section_header(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["purple_dark"])

    # Vertical accent bar
    add_rect(slide, Inches(0), Inches(2.5), Inches(0.4), Inches(2.5),
             fill=BRAND["purple_primary"])

    add_textbox(slide, Inches(0.9), Inches(2.4), Inches(11.5), Inches(0.5),
                data["kicker"].upper(), size=12, bold=True, color=BRAND["purple_light"])
    add_textbox(slide, Inches(0.9), Inches(2.95), Inches(11.5), Inches(1.6),
                data["title"], font_name=FONT_HEAD, size=48, bold=True,
                color=BRAND["bg_white"], line_spacing=1.1)

    # Preview bullets at bottom
    if "preview" in data:
        prev_top = Inches(5.0)
        prev = data["preview"]
        for i, p in enumerate(prev):
            y = prev_top + Inches(0.5) * i
            # Small accent
            add_rect(slide, Inches(0.9), y + Inches(0.18), Inches(0.18), Inches(0.05),
                     fill=BRAND["magenta"])
            add_textbox(slide, Inches(1.2), y, Inches(11), Inches(0.45),
                        p, font_name=FONT_BODY, size=15, color=BRAND["purple_light"])

    add_textbox(slide, Inches(0.4), Inches(7.12), Inches(12.53), Inches(0.3),
                f"{page_no} / {total}", size=8,
                color=BRAND["purple_light"], align=PP_ALIGN.RIGHT)

    set_notes(slide, f"Section divider for {data['title']}. {data['kicker']}.")


def build_stat_grid(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["bg_white"])
    add_brand_mark(slide)
    add_kicker(slide, "Findings")
    add_slide_title(slide, data["title"])
    add_key_message(slide, data["key_message"])

    stats = data["stats"]
    n = len(stats)
    grid_top = Inches(2.15)
    grid_h = Inches(2.55)
    margin_x = Inches(0.4)
    gap = Inches(0.18)

    cols = min(n, 4)
    rows = (n + cols - 1) // cols

    avail_w = int(SLIDE_W) - int(margin_x) * 2 - int(gap) * (cols - 1)
    tile_w = Emu(avail_w // cols)
    tile_h = Emu(int(grid_h) // rows) if rows > 1 else grid_h

    palette = [BRAND["purple_primary"], BRAND["purple_dark"],
               BRAND["purple_mid"], BRAND["magenta_dark"]]

    for idx, st in enumerate(stats):
        r = idx // cols
        c = idx % cols
        left = margin_x + Emu((int(tile_w) + int(gap)) * c)
        top = grid_top + Emu((int(tile_h) + int(gap)) * r)

        bg = palette[idx % len(palette)]
        add_rect(slide, left, top, tile_w, tile_h, fill=bg)

        # Big stat
        add_textbox(slide, left + Inches(0.25), top + Inches(0.2),
                    tile_w - Inches(0.5), Inches(0.95),
                    st["value"], font_name=FONT_HEAD, size=42, bold=True,
                    color=BRAND["bg_white"], line_spacing=1.0)

        # Thin rule
        add_rect(slide, left + Inches(0.25), top + Inches(1.1),
                 Inches(0.6), Emu(15875), fill=BRAND["purple_light"])

        # Label
        add_textbox(slide, left + Inches(0.25), top + Inches(1.2),
                    tile_w - Inches(0.5), Inches(0.7),
                    st["label"], font_name=FONT_BODY, size=11, bold=False,
                    color=BRAND["bg_white"], line_spacing=1.25)

        # Context
        if "context" in st:
            add_textbox(slide, left + Inches(0.25), top + tile_h - Inches(0.85),
                        tile_w - Inches(0.5), Inches(0.45),
                        st["context"], font_name=FONT_BODY, size=10,
                        italic=True, color=BRAND["purple_light"], line_spacing=1.25)

        # Source pill
        add_textbox(slide, left + Inches(0.25), top + tile_h - Inches(0.35),
                    tile_w - Inches(0.5), Inches(0.3),
                    st["source"], font_name=FONT_BODY, size=8,
                    color=BRAND["purple_light"])

    # Narrative band
    band_top = grid_top + Emu((int(tile_h) + int(gap)) * rows) + Inches(0.15)
    add_rect(slide, Inches(0.4), band_top, Inches(0.06), Inches(0.7),
             fill=BRAND["magenta"])
    add_textbox(slide, Inches(0.6), band_top, Inches(12.3), Inches(0.8),
                data["narrative"], font_name=FONT_BODY, size=12,
                color=BRAND["slate"], line_spacing=1.4, italic=True)

    add_footer(slide, page_no, total, deck_meta)
    set_notes(slide, "Walk through the four numbers in order. Anchor on +15% as the headline, then immediately reveal the dispersion — +36% bottom quintile vs ~0% top quintile. Close with the −40% attrition figure as a teaser for human-capital implications. All four numbers come directly from the paper's Abstract and Tables 2–3; no derived calculations.")


def build_distribution(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["bg_white"])
    add_brand_mark(slide)
    add_kicker(slide, data["kicker"])
    add_slide_title(slide, data["title"])
    add_key_message(slide, data["key_message"])

    bars = data["bars"]
    n = len(bars)
    chart_left = Inches(1.0)
    chart_top = Inches(2.5)
    chart_w = Inches(11.3)
    chart_h = Inches(3.4)

    # Baseline (zero line)
    baseline_y = chart_top + chart_h
    add_rect(slide, chart_left, baseline_y, chart_w, Emu(15875),
             fill=BRAND["slate"])

    # Each bar
    bar_gap = Inches(0.25)
    bar_total_w = int(chart_w) - int(bar_gap) * (n - 1)
    bar_w = Emu(bar_total_w // n)
    max_value = max(b["value"] for b in bars)
    palette = [BRAND["magenta_dark"], BRAND["magenta"], BRAND["purple_mid"],
               BRAND["berry"], BRAND["purple_light"]]

    for i, b in enumerate(bars):
        x = chart_left + Emu((int(bar_w) + int(bar_gap)) * i)
        # Bar height proportional to value
        bar_h_emu = int(int(chart_h) * (b["value"] / max(max_value, 1)))
        if bar_h_emu < int(Inches(0.15)):
            bar_h_emu = int(Inches(0.15))
        bar_h = Emu(bar_h_emu)
        bar_top = baseline_y - bar_h
        add_rect(slide, x, bar_top, bar_w, bar_h, fill=palette[i % len(palette)])

        # Value label on top
        add_textbox(slide, x, bar_top - Inches(0.5), bar_w, Inches(0.4),
                    b["display"], font_name=FONT_HEAD, size=18, bold=True,
                    color=BRAND["purple_dark"], align=PP_ALIGN.CENTER)

        # Quintile label below
        add_textbox(slide, x, baseline_y + Inches(0.1), bar_w, Inches(0.4),
                    b["label"], font_name=FONT_BODY, size=11, bold=True,
                    color=BRAND["slate"], align=PP_ALIGN.CENTER)

    # Annotations (left/right arrows)
    add_rect(slide, Inches(0.4), Inches(6.4), Inches(0.06), Inches(0.4),
             fill=BRAND["magenta"])
    add_textbox(slide, Inches(0.55), Inches(6.4), Inches(5.5), Inches(0.4),
                f"← {data['annotation_left']}", font_name=FONT_BODY, size=11,
                color=BRAND["magenta_dark"], bold=True)
    add_textbox(slide, Inches(7.0), Inches(6.4), Inches(5.7), Inches(0.4),
                f"{data['annotation_right']} →", font_name=FONT_BODY, size=11,
                color=BRAND["slate"], bold=True, align=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(0.4), Inches(6.85), Inches(12.53), Inches(0.25),
                f"Source: {data['source_note']}",
                size=8, color=BRAND["muted"], italic=True)

    add_footer(slide, page_no, total, deck_meta)
    set_notes(slide, "This bar chart visualizes what the headline number hides. The leftmost bar (Q1, lowest skill) shows +36% RPH; the rightmost (Q5, highest skill) shows essentially zero. Note that intermediate quintile values (Q2, Q3, Q4) are illustrative — the paper reports a broadly monotonic decline; the precise intermediate points are interpolated from Figures 4–5. The shape of the gradient is the message: gains concentrate at the bottom.")


def build_quote_pair(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["purple_tint"])
    add_brand_mark(slide)
    add_kicker(slide, data["kicker"])
    add_slide_title(slide, data["title"])

    # Primary quote
    pq_top = Inches(1.7)
    pq_h = Inches(2.0)
    add_rect(slide, Inches(0.4), pq_top, Inches(0.08), pq_h, fill=BRAND["purple_primary"])
    add_textbox(slide, Inches(0.7), pq_top, Inches(0.4), Inches(1.5),
                "“", font_name=FONT_HEAD, size=80, bold=True,
                color=BRAND["purple_primary"], line_spacing=0.9)
    add_textbox(slide, Inches(1.3), pq_top + Inches(0.15), Inches(11.5), Inches(1.4),
                data["primary_quote"], font_name=FONT_HEAD, size=20,
                color=BRAND["purple_dark"], line_spacing=1.35)
    add_textbox(slide, Inches(1.3), pq_top + Inches(1.55), Inches(11.5), Inches(0.35),
                f"— {data['primary_attrib']}", font_name=FONT_BODY, size=11, bold=True,
                color=BRAND["slate"])

    # Secondary quote
    sq_top = Inches(3.85)
    sq_h = Inches(1.6)
    add_rect(slide, Inches(0.4), sq_top, Inches(12.53), sq_h,
             fill=BRAND["bg_white"], line=BRAND["rule"])
    add_rect(slide, Inches(0.4), sq_top, Inches(0.08), sq_h, fill=BRAND["magenta"])
    add_textbox(slide, Inches(0.7), sq_top + Inches(0.2), Inches(12), Inches(0.9),
                data["secondary_quote"], font_name=FONT_HEAD, size=16,
                color=BRAND["ink"], line_spacing=1.35, italic=True)
    add_textbox(slide, Inches(0.7), sq_top + sq_h - Inches(0.45), Inches(12), Inches(0.35),
                f"— {data['secondary_attrib']}", font_name=FONT_BODY, size=10, bold=True,
                color=BRAND["muted"])

    # Implication band
    imp_top = Inches(5.65)
    add_rect(slide, Inches(0.4), imp_top, Inches(12.53), Inches(1.3),
             fill=BRAND["purple_dark"])
    add_textbox(slide, Inches(0.7), imp_top + Inches(0.2), Inches(11.9), Inches(0.4),
                data["implication_label"].upper(), font_name=FONT_HEAD, size=11, bold=True,
                color=BRAND["purple_light"])
    add_textbox(slide, Inches(0.7), imp_top + Inches(0.55), Inches(11.9), Inches(0.7),
                data["implication"], font_name=FONT_BODY, size=14,
                color=BRAND["bg_white"], line_spacing=1.3)

    add_footer(slide, page_no, total, deck_meta)
    set_notes(slide, "Two direct quotes from the paper, paired to set up the thesis. The first establishes skill-equalization. The second quantifies it concretely (2-month new hires match 6+ month veterans). The dark implication band makes the historical contrast explicit — this is the first IT wave that's not skill-biased.")


def build_three_column_rich(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["bg_white"])
    add_brand_mark(slide)
    add_kicker(slide, data["kicker"])
    add_slide_title(slide, data["title"])
    add_key_message(slide, data["key_message"])

    cols_data = data["columns"]
    col_w = Inches(4.0)
    col_h = Inches(4.4)
    col_top = Inches(2.2)
    gap = Inches(0.18)
    start_x = Inches(0.4)
    palette = [BRAND["purple_primary"], BRAND["purple_mid"], BRAND["berry"]]

    for i, col in enumerate(cols_data):
        x = start_x + Emu((int(col_w) + int(gap)) * i)
        # Big icon header
        add_rect(slide, x, col_top, col_w, Inches(0.85), fill=palette[i])
        add_textbox(slide, x + Inches(0.25), col_top + Inches(0.05),
                    Inches(0.7), Inches(0.7),
                    col["icon_text"], font_name=FONT_HEAD, size=32, bold=True,
                    color=BRAND["bg_white"])
        add_textbox(slide, x + Inches(1.0), col_top + Inches(0.2),
                    col_w - Inches(1.2), Inches(0.55),
                    col["header"], font_name=FONT_HEAD, size=14, bold=True,
                    color=BRAND["bg_white"], line_spacing=1.15)

        # Body card
        body_top = col_top + Inches(0.85)
        body_h = col_h - Inches(0.85)
        add_rect(slide, x, body_top, col_w, body_h,
                 fill=BRAND["bg_off"], line=BRAND["rule"])

        # Body text
        add_textbox(slide, x + Inches(0.25), body_top + Inches(0.2),
                    col_w - Inches(0.5), Inches(1.4),
                    col["body"], font_name=FONT_BODY, size=13,
                    color=BRAND["ink"], line_spacing=1.4)

        # Evidence sub-block
        ev_top = body_top + Inches(1.7)
        ev_h = Inches(1.4)
        add_rect(slide, x + Inches(0.15), ev_top,
                 col_w - Inches(0.3), ev_h,
                 fill=BRAND["evidence_bg"])
        add_rect(slide, x + Inches(0.15), ev_top,
                 Inches(0.05), ev_h, fill=BRAND["purple_primary"])
        add_textbox(slide, x + Inches(0.35), ev_top + Inches(0.1),
                    col_w - Inches(0.5), Inches(0.3),
                    col["evidence_label"].upper(),
                    font_name=FONT_HEAD, size=9, bold=True,
                    color=BRAND["purple_primary"])
        add_textbox(slide, x + Inches(0.35), ev_top + Inches(0.4),
                    col_w - Inches(0.5), ev_h - Inches(0.5),
                    col["evidence"], font_name=FONT_BODY, size=11,
                    color=BRAND["purple_dark"], line_spacing=1.35, italic=True)

        # Source
        add_textbox(slide, x + Inches(0.25), body_top + body_h - Inches(0.3),
                    col_w - Inches(0.5), Inches(0.25),
                    col["source"], font_name=FONT_BODY, size=8,
                    color=BRAND["muted"])

    add_footer(slide, page_no, total, deck_meta)
    set_notes(slide, "Three mechanisms is the analytical heart of the paper. The outage natural experiment (column 1) is the most clever — rules out 'AI just helps in real time'. Column 2 (language uplift) makes this directly relevant to global outsourcing. Column 3 (rare-problem coverage) is counter-intuitive but logically consistent: less training data, but lower human baseline.")


def build_two_column_rich(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["bg_white"])
    add_brand_mark(slide)
    add_kicker(slide, data["kicker"])
    add_slide_title(slide, data["title"])
    add_key_message(slide, data["key_message"])

    col_w = Inches(5.95)
    col_h = Inches(3.0)
    col_top = Inches(2.15)

    # Left
    add_rect(slide, Inches(0.4), col_top, col_w, Inches(0.5), fill=BRAND["slate"])
    add_textbox(slide, Inches(0.6), col_top + Inches(0.1),
                col_w - Inches(0.4), Inches(0.4),
                data["left"]["header"].upper(), font_name=FONT_HEAD, size=12, bold=True,
                color=BRAND["bg_white"])
    add_rect(slide, Inches(0.4), col_top + Inches(0.5), col_w, col_h - Inches(0.5),
             fill=BRAND["bg_off"], line=BRAND["rule"])
    add_paragraphs(slide, Inches(0.7), col_top + Inches(0.7),
                   col_w - Inches(0.5), col_h - Inches(0.8),
                   data["left"]["points"], size=14, color=BRAND["ink"],
                   bullet=True, line_spacing=1.5, space_after=4)

    # Right
    add_rect(slide, Inches(6.95), col_top, col_w, Inches(0.5), fill=BRAND["purple_primary"])
    add_textbox(slide, Inches(7.15), col_top + Inches(0.1),
                col_w - Inches(0.4), Inches(0.4),
                data["right"]["header"].upper(), font_name=FONT_HEAD, size=12, bold=True,
                color=BRAND["bg_white"])
    add_rect(slide, Inches(6.95), col_top + Inches(0.5), col_w, col_h - Inches(0.5),
             fill=BRAND["purple_tint"], line=BRAND["rule"])
    add_paragraphs(slide, Inches(7.25), col_top + Inches(0.7),
                   col_w - Inches(0.5), col_h - Inches(0.8),
                   data["right"]["points"], size=14, color=BRAND["purple_dark"],
                   bullet=True, line_spacing=1.5, space_after=4)

    # Callout block
    co = data["callout"]
    co_top = col_top + col_h + Inches(0.25)
    co_h = Inches(1.2)
    add_rect(slide, Inches(0.4), co_top, Inches(12.53), co_h, fill=BRAND["purple_dark"])
    add_rect(slide, Inches(0.4), co_top, Inches(0.1), co_h, fill=BRAND["magenta"])
    add_textbox(slide, Inches(0.7), co_top + Inches(0.2), Inches(12), Inches(0.4),
                co["label"].upper(), font_name=FONT_HEAD, size=11, bold=True,
                color=BRAND["purple_light"])
    add_textbox(slide, Inches(0.7), co_top + Inches(0.55), Inches(12), Inches(0.6),
                co["text"], font_name=FONT_BODY, size=14,
                color=BRAND["bg_white"], line_spacing=1.3)

    add_textbox(slide, Inches(0.4), Inches(6.85), Inches(12.53), Inches(0.25),
                f"Source: {data['source_note']}",
                size=8, color=BRAND["muted"], italic=True)

    add_footer(slide, page_no, total, deck_meta)
    set_notes(slide, f"Theme: {data['title']}. Use the contrast between columns to drive the narrative — left is baseline, right is AI-enabled. The dark callout band at the bottom is the strategic 'so what'.")


def build_insight_card_rich(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["bg_off"])
    add_brand_mark(slide)
    add_kicker(slide, data["kicker"])
    add_slide_title(slide, data["title"])
    add_key_message(slide, data["key_message"])

    cards = data["cards"]
    card_w = Inches(5.95)
    card_h = Inches(4.0)
    card_top = Inches(2.15)
    risk_colors = [BRAND["magenta"], BRAND["purple_dark"]]

    for i, card in enumerate(cards):
        x = Inches(0.4) + Emu((int(card_w) + int(Inches(0.2))) * i)
        # Tag pill
        add_rect(slide, x, card_top, card_w, Inches(0.5), fill=risk_colors[i])
        add_textbox(slide, x + Inches(0.25), card_top + Inches(0.1),
                    Inches(1.5), Inches(0.35),
                    card["tag"].upper(), font_name=FONT_HEAD, size=11, bold=True,
                    color=BRAND["bg_white"])
        add_textbox(slide, x + Inches(1.5), card_top + Inches(0.05),
                    card_w - Inches(1.7), Inches(0.45),
                    card["header"], font_name=FONT_HEAD, size=18, bold=True,
                    color=BRAND["bg_white"])

        # Body card
        body_top = card_top + Inches(0.5)
        body_h = card_h - Inches(0.5)
        add_rect(slide, x, body_top, card_w, body_h,
                 fill=BRAND["bg_white"], line=BRAND["rule"])

        # Body text
        add_textbox(slide, x + Inches(0.3), body_top + Inches(0.2),
                    card_w - Inches(0.6), Inches(1.3),
                    card["body"], font_name=FONT_BODY, size=14,
                    color=BRAND["ink"], line_spacing=1.4)

        # Evidence sub-block
        ev_top = body_top + Inches(1.65)
        ev_h = Inches(1.4)
        add_rect(slide, x + Inches(0.2), ev_top,
                 card_w - Inches(0.4), ev_h,
                 fill=BRAND["evidence_bg"])
        add_rect(slide, x + Inches(0.2), ev_top,
                 Inches(0.05), ev_h, fill=risk_colors[i])
        add_textbox(slide, x + Inches(0.4), ev_top + Inches(0.1),
                    card_w - Inches(0.6), Inches(0.3),
                    card["evidence_label"].upper(),
                    font_name=FONT_HEAD, size=9, bold=True,
                    color=risk_colors[i])
        add_textbox(slide, x + Inches(0.4), ev_top + Inches(0.4),
                    card_w - Inches(0.6), ev_h - Inches(0.5),
                    card["evidence"], font_name=FONT_BODY, size=11,
                    color=BRAND["slate"], line_spacing=1.35, italic=True)

        # Implication
        add_textbox(slide, x + Inches(0.3), body_top + body_h - Inches(0.55),
                    card_w - Inches(0.6), Inches(0.45),
                    card["implication"], font_name=FONT_BODY, size=12, bold=True,
                    color=risk_colors[i], line_spacing=1.3)

    add_textbox(slide, Inches(0.4), Inches(6.4), Inches(12.53), Inches(0.25),
                f"Source: {data['source_note']}",
                size=8, color=BRAND["muted"], italic=True)

    add_footer(slide, page_no, total, deck_meta)
    set_notes(slide, "This is where the deck transitions from 'what happened' to 'what to worry about'. Both risks come from §8 of the paper. Risk 1 (top-talent flattening) is supported by the small quality-decline at the top. Risk 2 (training data decay) is the policy-level question the paper raises but cannot resolve. Frame these as 'today's win, tomorrow's structural problem'.")


def build_action_table(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["bg_white"])
    add_brand_mark(slide)
    add_kicker(slide, data["kicker"])
    add_slide_title(slide, data["title"])
    add_key_message(slide, data["key_message"])

    actions = data["actions"]
    table_top = Inches(2.15)
    row_h = Inches(0.85)
    n_rows = len(actions) + 1  # header row

    # Column widths
    col_num_w = Inches(0.7)
    col_action_w = Inches(3.4)
    col_what_w = Inches(4.5)
    col_outcome_w = Inches(3.93)
    margin_x = Inches(0.4)

    # Header row
    add_rect(slide, margin_x, table_top, Inches(12.53), row_h,
             fill=BRAND["purple_dark"])
    headers = [("#", col_num_w, PP_ALIGN.CENTER),
               ("Action", col_action_w, PP_ALIGN.LEFT),
               ("What to do", col_what_w, PP_ALIGN.LEFT),
               ("Expected outcome", col_outcome_w, PP_ALIGN.LEFT)]
    x_offset = margin_x
    for label, w, align in headers:
        add_textbox(slide, x_offset + Inches(0.15), table_top + Inches(0.27),
                    w - Inches(0.3), Inches(0.4),
                    label.upper(), font_name=FONT_HEAD, size=11, bold=True,
                    color=BRAND["purple_light"], align=align)
        x_offset = x_offset + w

    # Body rows
    for i, act in enumerate(actions):
        y = table_top + row_h * (i + 1)
        bg = BRAND["bg_off"] if i % 2 == 0 else BRAND["bg_white"]
        add_rect(slide, margin_x, y, Inches(12.53), row_h,
                 fill=bg, line=BRAND["soft_rule"])

        x = margin_x
        # Number
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         x + Inches(0.15), y + Inches(0.15),
                                         Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BRAND["purple_primary"]
        circle.line.fill.background()
        add_textbox(slide, x + Inches(0.15), y + Inches(0.22),
                    Inches(0.5), Inches(0.4),
                    f"{i+1}", size=14, bold=True,
                    color=BRAND["bg_white"], align=PP_ALIGN.CENTER, font_name=FONT_HEAD)
        x = x + col_num_w

        # Action header
        add_textbox(slide, x + Inches(0.15), y + Inches(0.22),
                    col_action_w - Inches(0.3), Inches(0.5),
                    act["header"], font_name=FONT_HEAD, size=14, bold=True,
                    color=BRAND["purple_dark"], line_spacing=1.2)
        x = x + col_action_w

        # What
        add_textbox(slide, x + Inches(0.15), y + Inches(0.18),
                    col_what_w - Inches(0.3), Inches(0.6),
                    act["what"], font_name=FONT_BODY, size=11,
                    color=BRAND["ink"], line_spacing=1.3)
        x = x + col_what_w

        # Outcome
        add_textbox(slide, x + Inches(0.15), y + Inches(0.18),
                    col_outcome_w - Inches(0.3), Inches(0.6),
                    act["outcome"], font_name=FONT_BODY, size=11, italic=True,
                    color=BRAND["slate"], line_spacing=1.3)

    add_footer(slide, page_no, total, deck_meta)
    set_notes(slide, "Five concrete moves derived from the paper's findings. Action #2 (new-hire retention as a hard KPI) is the single highest-leverage move — the −40% attrition figure has more durable ROI than the +15% productivity number. Action #5 (protect top-talent originality) directly addresses the long-term risk discussed two slides ago. Walk through the table column by column — What to do, then Expected outcome.")


def build_summary(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["purple_dark"])

    # Top accent
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08),
             fill=BRAND["purple_primary"])

    add_textbox(slide, Inches(0.4), Inches(0.5), Inches(8), Inches(0.3),
                data["kicker"].upper(), size=11, bold=True,
                color=BRAND["purple_light"])
    add_textbox(slide, Inches(0.4), Inches(0.85), Inches(12.53), Inches(0.7),
                data["title"], font_name=FONT_HEAD, size=32, bold=True,
                color=BRAND["bg_white"])
    add_textbox(slide, Inches(0.4), Inches(1.55), Inches(12.53), Inches(0.6),
                data["key_message"], font_name=FONT_HEAD, size=15, italic=True,
                color=BRAND["purple_light"], line_spacing=1.3)

    # Three pillars
    pillars = data["pillars"]
    col_w = Inches(4.0)
    col_h = Inches(3.4)
    col_top = Inches(2.5)
    gap = Inches(0.18)
    palette = [BRAND["purple_primary"], BRAND["berry"], BRAND["magenta"]]

    for i, p in enumerate(pillars):
        x = Inches(0.4) + Emu((int(col_w) + int(gap)) * i)
        # Card
        add_rect(slide, x, col_top, col_w, col_h, fill=BRAND["bg_white"])
        # Top accent
        add_rect(slide, x, col_top, col_w, Inches(0.15), fill=palette[i])
        # Big number
        add_textbox(slide, x + Inches(0.3), col_top + Inches(0.4),
                    Inches(1.5), Inches(0.7),
                    f"0{i+1}", font_name=FONT_HEAD, size=36, bold=True,
                    color=palette[i], line_spacing=1.0)
        # Header
        add_textbox(slide, x + Inches(0.3), col_top + Inches(1.15),
                    col_w - Inches(0.6), Inches(0.5),
                    p["header"], font_name=FONT_HEAD, size=16, bold=True,
                    color=BRAND["purple_dark"])
        # Body
        add_textbox(slide, x + Inches(0.3), col_top + Inches(1.7),
                    col_w - Inches(0.6), col_h - Inches(1.85),
                    p["body"], font_name=FONT_BODY, size=12,
                    color=BRAND["ink"], line_spacing=1.4)

    # Closing line
    cl_top = col_top + col_h + Inches(0.25)
    add_rect(slide, Inches(0.4), cl_top, Inches(0.08), Inches(0.7),
             fill=BRAND["magenta"])
    add_textbox(slide, Inches(0.7), cl_top, Inches(12.2), Inches(0.7),
                data["closing_line"], font_name=FONT_HEAD, size=15, italic=True,
                color=BRAND["purple_light"], line_spacing=1.4)

    # Footer (light on dark)
    add_textbox(slide, Inches(0.4), Inches(7.12), Inches(12.53), Inches(0.3),
                f"{page_no} / {total}", size=8,
                color=BRAND["purple_light"], align=PP_ALIGN.RIGHT)

    set_notes(slide, "Synthesis slide — close out the substantive content with three takeaways: what changed, what it means, what to watch. The closing line emphasizes urgency: redesign HR, ROI, and data-attribution mechanisms now while gains accrue.")


def build_references(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["bg_off"])
    add_brand_mark(slide)
    add_kicker(slide, data["kicker"])
    add_slide_title(slide, data["title"])

    # Three sections side by side
    sec_w = Inches(4.0)
    sec_h = Inches(4.5)
    sec_top = Inches(1.7)
    gap = Inches(0.18)

    sections = [data["primary"], data["method"], data["citations"]]
    palette = [BRAND["purple_primary"], BRAND["purple_mid"], BRAND["berry"]]

    for i, sec in enumerate(sections):
        x = Inches(0.4) + Emu((int(sec_w) + int(gap)) * i)
        # Header
        add_rect(slide, x, sec_top, sec_w, Inches(0.5), fill=palette[i])
        add_textbox(slide, x + Inches(0.2), sec_top + Inches(0.1),
                    sec_w - Inches(0.4), Inches(0.4),
                    sec["header"].upper(), font_name=FONT_HEAD, size=11, bold=True,
                    color=BRAND["bg_white"])
        # Body
        add_rect(slide, x, sec_top + Inches(0.5), sec_w, sec_h - Inches(0.5),
                 fill=BRAND["bg_white"], line=BRAND["rule"])
        add_paragraphs(slide, x + Inches(0.2), sec_top + Inches(0.7),
                       sec_w - Inches(0.4), sec_h - Inches(0.85),
                       sec["items"], size=10, color=BRAND["ink"],
                       bullet=True, line_spacing=1.4, space_after=4)

    add_footer(slide, page_no, total, deck_meta)
    set_notes(slide, "All claims trace back to the cited paper. Use this slide if asked about methodology — the two-way fixed effects DiD design and the outage-based natural experiment are what give the findings their causal weight.")


def build_closing(prs, data, deck_meta, page_no, total):
    slide = add_blank_slide(prs)
    set_bg(slide, BRAND["purple_dark"])

    add_rect(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(0.05),
             fill=BRAND["purple_primary"])

    add_textbox(slide, Inches(0.4), Inches(2.5), Inches(12.53), Inches(1),
                data["title"], font_name=FONT_HEAD, size=58, bold=True,
                color=BRAND["bg_white"], align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.4), Inches(3.6), Inches(12.53), Inches(0.6),
                data["subtitle"], font_name=FONT_BODY, size=18,
                color=BRAND["purple_light"], align=PP_ALIGN.CENTER)

    if "qr_text" in data:
        add_textbox(slide, Inches(0.4), Inches(5.0), Inches(12.53), Inches(0.4),
                    data["qr_text"], font_name=FONT_HEAD, size=14, bold=True,
                    color=BRAND["berry"], align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(0.4), Inches(5.4), Inches(12.53), Inches(0.4),
                    data["url_text"], font_name=FONT_BODY, size=12,
                    color=BRAND["purple_light"], align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.4), Inches(6.8), Inches(12.53), Inches(0.4),
                data["footer_meta"], font_name=FONT_BODY, size=10,
                color=BRAND["berry"], align=PP_ALIGN.CENTER)
    set_notes(slide, "Close with an invitation to dialogue. Remind the audience that this is one paper in the broader Stanford Enterprise AI Playbook lineage. Happy to dive deeper on any specific finding.")


# ---------- Dispatcher ----------

BUILDERS = {
    "cover": build_cover,
    "context": build_context,
    "agenda": build_agenda,
    "section_header": build_section_header,
    "stat_grid": build_stat_grid,
    "distribution": build_distribution,
    "quote_pair": build_quote_pair,
    "three_column_rich": build_three_column_rich,
    "two_column_rich": build_two_column_rich,
    "insight_card_rich": build_insight_card_rich,
    "action_table": build_action_table,
    "summary": build_summary,
    "references": build_references,
    "closing": build_closing,
}


def generate(outline_path: Path, output_path: Path):
    with outline_path.open() as f:
        outline = json.load(f)

    deck_meta = outline["deck_meta"]
    slides = outline["slides"]
    total = len(slides)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for idx, slide_data in enumerate(slides):
        page_no = idx + 1
        builder = BUILDERS.get(slide_data["type"])
        if not builder:
            print(f"WARN: no builder for type {slide_data['type']}")
            continue
        builder(prs, slide_data, deck_meta, page_no, total)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"OK: wrote {output_path} ({total} slides)")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python3 generate_deck.py <outline.json> <output.pptx>")
        sys.exit(1)
    generate(Path(sys.argv[1]), Path(sys.argv[2]))
