"""Low-level drawing helpers for Anthropic brand.

All text routes through these so font, color, palette stay consistent.
"""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from . import style as S


def add_rect(slide, x, y, w, h, *,
             fill=None, line=None, line_width=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width)
    # strip effects defensively
    sp = shp._element
    try:
        spPr = sp.find(qn("p:spPr"))
        if spPr is not None:
            for tag in ("a:effectLst", "a:effectDag"):
                for el in spPr.findall(qn(tag)):
                    spPr.remove(el)
    except Exception:
        pass
    return shp


def add_textbox(slide, x, y, w, h, text, *,
                font=None, font_size=Pt(11), bold=False, italic=False,
                color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                char_spacing=0, line_spacing=None, margin=Inches(0)):
    if color is None:
        color = S.DARK
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = margin
    tf.margin_bottom = margin
    if anchor is not None:
        tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    # multi-line text support
    lines = text.split("\n") if "\n" in text else [text]
    first = True
    for line in lines:
        if first:
            para = p
            first = False
        else:
            para = tf.add_paragraph()
            para.alignment = align
            if line_spacing is not None:
                para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        f = run.font
        f.name = font or S.FONT_BODY
        f.size = font_size
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
        if char_spacing:
            rPr = run._r.get_or_add_rPr()
            rPr.set("spc", str(int(char_spacing * 100)))
    return tb


def add_paragraphs(slide, x, y, w, h, runs, *,
                   anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
                   line_spacing=1.25, space_after=Pt(2),
                   margin=Inches(0.04)):
    """Add a textbox with structured paragraphs.

    Each item in `runs` is a paragraph spec. A spec can be:
      - {text, size, bold, italic, color, font, align, line_spacing, space_after}
      - {"runs": [run_spec, ...], align, line_spacing, space_after}
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = margin
    tf.margin_bottom = margin
    if anchor is not None:
        tf.vertical_anchor = anchor

    first = True
    for spec in runs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = spec.get("align", align)
        p.line_spacing = spec.get("line_spacing", line_spacing)
        p.space_after = spec.get("space_after", space_after)

        run_specs = spec["runs"] if "runs" in spec else [spec]
        for rs in run_specs:
            run = p.add_run()
            run.text = rs.get("text", "")
            f = run.font
            f.name = rs.get("font", S.FONT_BODY)
            f.size = rs.get("size", Pt(11))
            f.bold = rs.get("bold", False)
            f.italic = rs.get("italic", False)
            f.color.rgb = rs.get("color", S.DARK)
            cs = rs.get("char_spacing", 0)
            if cs:
                rPr = run._r.get_or_add_rPr()
                rPr.set("spc", str(int(cs * 100)))
    return tb


def add_hrule(slide, x, y, w, *, color=None, height=Pt(1.2)):
    """Add a thin horizontal rule."""
    if color is None:
        color = S.DARK
    add_rect(slide, x, y, w, height, fill=color, line=color)


def fill_background(slide, color):
    """Fill entire slide with a solid color."""
    add_rect(slide, Inches(0), Inches(0), S.SLIDE_W, S.SLIDE_H,
             fill=color, line=color)
