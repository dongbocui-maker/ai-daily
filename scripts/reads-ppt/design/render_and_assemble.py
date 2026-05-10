#!/usr/bin/env python3
"""
Two-step approach:
1. Generate one HTML file PER slide (to control page breaks)
2. Render each to PDF via weasyprint
3. Convert PDF pages to PNG via pdftoppm
4. Assemble into PPTX
"""
import json, sys, re, subprocess, glob, os
from pathlib import Path
from weasyprint import HTML
from pptx import Presentation
from pptx.util import Inches, Emu

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PROJ = Path(__file__).parent.parent

def main():
    html_path = PROJ / "output" / "creative-deck.html"
    slide_dir = PROJ / "output" / "slides-html"
    png_dir = PROJ / "output" / "slides-png"
    pdf_dir = PROJ / "output" / "slides-pdf"

    for d in [slide_dir, png_dir, pdf_dir]:
        d.mkdir(parents=True, exist_ok=True)

    with html_path.open() as f:
        content = f.read()

    # Extract CSS (everything between <style> and </style>)
    css_match = re.search(r'<style>(.*?)</style>', content, re.DOTALL)
    css_part = css_match.group(0) if css_match else ""

    # Split into slide blocks
    slide_pattern = re.compile(r'(<div class="slide .*?</div>\s*</div>)', re.DOTALL)
    slide_blocks = slide_pattern.findall(content)

    print(f"Found {len(slide_blocks)} slides")

    for i, block in enumerate(slide_blocks):
        pn = i + 1

        # Build standalone HTML for this slide
        full_html = f"""<!DOCTYPE html>
<html><head><meta charset='utf-8'>
{css_part}
<style>
  @page {{ size: 1920px 1080px; margin: 0; }}
  body {{ margin: 0; background: #0a0a0a; }}
</style>
</head><body>
{block}
</body></html>"""

        html_file = slide_dir / f"slide-{pn:02d}.html"
        html_file.write_text(full_html)

        pdf_file = pdf_dir / f"slide-{pn:02d}.pdf"
        try:
            HTML(string=full_html).write_pdf(str(pdf_file))
        except Exception as e:
            print(f"  [{pn:02d}] PDF FAILED: {e}")
            continue

        # Convert PDF to PNG
        png_prefix = png_dir / f"slide-{pn:02d}"
        result = subprocess.run(
            ["pdftoppm", "-png", "-r", "150", "-singlefile",
             str(pdf_file), str(png_prefix)],
            capture_output=True, text=True
        )
        actual_png = png_dir / f"slide-{pn:02d}.png"
        if actual_png.exists():
            print(f"  [{pn:02d}] OK ({actual_png.stat().st_size:,} bytes)")
        else:
            print(f"  [{pn:02d}] PNG FAILED: {result.stderr[:200]}")

    # Assemble PPTX
    print("\nAssembling PPTX...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    pngs = sorted(glob.glob(str(png_dir / "slide-*.png")))
    for png_path in pngs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(png_path, Inches(0), Inches(0), SLIDE_W, SLIDE_H)

    pptx_path = PROJ / "output" / "brynjolfsson-creative-v3.pptx"
    prs.save(str(pptx_path))
    print(f"OK: wrote {pptx_path} ({len(pngs)} slides)")

if __name__ == "__main__":
    main()
