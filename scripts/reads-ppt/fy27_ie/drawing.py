"""Low-level drawing primitives — every text/shape goes through these
so we keep consistent font, color, and zero shadow/gradient."""
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from . import style as S


def _strip_effects(shape):
    """Remove shadow / glow effects defensively (python-pptx defaults to none,
    but layout templates can carry them). Also disable line if not requested."""
    sp = shape.shape_elem if hasattr(shape, "shape_elem") else shape._element
    # python-pptx Shape has _element
    try:
        spPr = sp.find(qn("p:spPr"))
        if spPr is not None:
            for tag in ("a:effectLst", "a:effectDag"):
                for el in spPr.findall(qn(tag)):
                    spPr.remove(el)
    except Exception:
        pass


def add_rect(slide, x, y, w, h, fill=None, line=None, line_width=0.5):
    """Add a filled rectangle. fill / line are RGBColor or None."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_width)
    _strip_effects(shp)
    return shp


def add_textbox(slide, x, y, w, h, text, *,
                font_size=Pt(10), bold=False, italic=False,
                color=S.INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                char_spacing=0, font=S.FONT, line_spacing=None,
                margin=Inches(0)):
    """Add a textbox with FY27 defaults (Calibri, palette color)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = margin
    tf.margin_bottom = margin
    if anchor is not None:
        tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = font_size
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if char_spacing:
        # spc is in 1/100 pt
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(char_spacing * 100)))
    return tb


def add_paragraphs(slide, x, y, w, h, runs, *,
                   anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
                   line_spacing=1.15, space_after=Pt(2),
                   margin=Inches(0.04)):
    """Add a textbox with multiple paragraphs.

    `runs` is a list of dicts: each dict can be a paragraph spec OR a list of run specs.

    Paragraph spec (single run): {text, size, bold, italic, color, align, bullet}
    Multi-run paragraph: {"runs": [run_spec, run_spec, ...], align, bullet}
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = margin
    tf.margin_bottom = margin
    if anchor is not None:
        tf.vertical_anchor = anchor

    first = True
    for spec in runs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = spec.get("align", align)
        p.line_spacing = spec.get("line_spacing", line_spacing)
        p.space_after = spec.get("space_after", space_after)

        if "runs" in spec:
            run_specs = spec["runs"]
        else:
            run_specs = [spec]

        for rs in run_specs:
            run = p.add_run()
            t = rs.get("text", "")
            run.text = t
            f = run.font
            f.name = rs.get("font", S.FONT)
            f.size = rs.get("size", Pt(10))
            f.bold = rs.get("bold", False)
            f.italic = rs.get("italic", False)
            color = rs.get("color", S.INK)
            f.color.rgb = color
            cs = rs.get("char_spacing", 0)
            if cs:
                rPr = run._r.get_or_add_rPr()
                rPr.set("spc", str(int(cs * 100)))
    return tb


def add_card(slide, x, y, w, h, *,
             fill=S.WHITE, line=S.BORDER, line_width=0.5,
             accent_top=False, accent_color=None, accent_h=Inches(0.04)):
    """White card with optional top purple accent bar."""
    add_rect(slide, x, y, w, h, fill=fill, line=line, line_width=line_width)
    if accent_top:
        add_rect(slide, x, y, w, accent_h,
                 fill=accent_color or S.PURPLE,
                 line=accent_color or S.PURPLE)


def add_pill(slide, x, y, h, text, *,
             pad_x=Inches(0.20), fill=S.PURPLE, color=S.WHITE,
             size=S.SZ_PILL, char_spacing=S.CSP_PILL):
    """Pill button — height fixed, width auto from text length estimate."""
    # crude width estimate: ~0.085" per char at 9pt bold uppercase
    text_w = Inches(0.085 * len(text)) + pad_x * 2
    add_rect(slide, x, y, text_w, h, fill=fill, line=fill)
    add_textbox(slide, x, y, text_w, h, text,
                font_size=size, bold=True, color=color,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                char_spacing=char_spacing, margin=Inches(0))
    return text_w


def add_bullet_line(specs):
    """Helper: build a paragraph spec for a bullet line with ● prefix."""
    bullet_run = {"text": S.BULLET_CHAR + "  ", "size": S.SZ_BULLET,
                  "color": S.PURPLE, "bold": True}
    return {"runs": [bullet_run] + specs, "line_spacing": 1.18, "space_after": Pt(3)}
